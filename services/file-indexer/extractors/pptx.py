"""PPTX extractor using python-pptx (WARP-435 / ADR-003 Phase 1).

PowerPoint .pptx files were not previously indexed. ADR-003 calls for
sentence-aware chunking + section-path enrichment, and PPTX is one of
the four "structured" formats whose layout naturally yields hierarchical
context (section name → slide title → body text).

``sectionPath`` rules:
  * ``[section_name, slide_title]`` when the slide belongs to a named
    PPTX *section* (the new-style sections that PowerPoint added in
    2010+, accessible via ``presentation.slides``' parent section
    metadata when present).
  * ``[slide_title]`` when the slide has a title but no enclosing
    section.
  * ``[]`` for untitled slides — caller falls back to ``[filename]``.

Body extraction walks every shape on every slide, concatenating any
``has_text_frame`` shape's text. Tables flatten cells with `" | "`.
Notes (speaker notes) are concatenated at the end of each slide's body.
"""
from __future__ import annotations

import logging
import os

from anchor_schema import NoneAnchor
from extractors.spans import Span
from extractors.types import ExtractedDoc

logger = logging.getLogger(__name__)


SUPPORTED_MIMES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }
)


def _slide_title(slide) -> str:
    """Return the slide's title text or an empty string."""
    try:
        title_shape = slide.shapes.title
    except Exception:
        return ""
    if title_shape is None:
        return ""
    try:
        return (title_shape.text or "").strip()
    except Exception:
        return ""


def _shape_text(shape) -> str:
    """Extract text from a single shape (text frame, table, group)."""
    out: list[str] = []
    if getattr(shape, "has_text_frame", False):
        try:
            t = shape.text_frame.text.strip()
            if t:
                out.append(t)
        except Exception:
            pass
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    out.append(row_text)
        except Exception:
            pass
    # Group shapes — recurse.
    if hasattr(shape, "shapes"):
        try:
            for child in shape.shapes:
                t = _shape_text(child)
                if t:
                    out.append(t)
        except Exception:
            pass
    return "\n".join(out)


def _section_name_for_slide(presentation, slide) -> str:
    """Best-effort lookup of the section name containing ``slide``.

    PPTX's section metadata lives in ``presentation.element`` XML under
    the ``p:sectionLst`` extension; python-pptx exposes it via
    ``presentation.slides`` only on newer versions (>=0.6.22). We
    introspect defensively and return ``""`` when the API isn't
    available — section info is nice-to-have, not load-bearing.
    """
    try:
        sections = getattr(presentation, "slide_layouts", None)
        # python-pptx >= 0.6.22 exposes presentation.sections;
        # gate on hasattr to avoid AttributeError on older builds.
        pres_sections = getattr(presentation, "sections", None)
        if pres_sections is None:
            return ""
        for section in pres_sections:
            try:
                if slide in section.slides:
                    return (section.name or "").strip()
            except Exception:
                continue
    except Exception:
        pass
    return ""


def extract(path: str) -> ExtractedDoc:
    # Local import — python-pptx pulls in lxml at module load time.
    from pptx import Presentation  # noqa: PLC0415

    presentation = Presentation(path)
    filename = os.path.basename(path) or "presentation"

    # WARP-287 + WARP-435: emit one Span per slide. PPTX is non-MVP for
    # positional anchors (no slide-anchor kind in the schema), so every
    # Span carries ``NoneAnchor``; the per-slide section path (section
    # name → slide title) rides on ``Span.section_path``, falling back to
    # ``[filename]`` for untitled / unsectioned slides.
    spans: list[Span] = []
    word_count = 0

    for slide in presentation.slides:
        title = _slide_title(slide)
        section_name = _section_name_for_slide(presentation, slide)

        if section_name and title:
            path_list = [section_name, title]
        elif title:
            path_list = [title]
        elif section_name:
            path_list = [section_name]
        else:
            path_list = [filename]

        slide_parts: list[str] = []
        if title:
            slide_parts.append(f"# {title}")
        for shape in slide.shapes:
            if shape is getattr(slide.shapes, "title", None):
                continue  # already emitted as the title line
            t = _shape_text(shape)
            if t:
                slide_parts.append(t)

        # Speaker notes (often substantive — index them).
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[notes] {notes}")
        except Exception:
            pass

        if not slide_parts:
            continue

        slide_text = "\n".join(slide_parts)
        word_count += len(slide_text.split())
        spans.append(
            Span(text=slide_text, anchor=NoneAnchor(), section_path=path_list)
        )

    return ExtractedDoc(
        spans=spans,
        language=None,
        metadata={
            "extractor_name": "pptx",
            "extractor_version": "2",
            "slide_count": len(presentation.slides),
            "word_count": word_count,
        },
        warnings=[],
    )
